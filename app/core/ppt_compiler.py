import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PPTCompiler:
    def __init__(self, repository):
        self.repo = repository
        
    def _hex_to_rgb(self, hex_str: str) -> RGBColor:
        """Converts a Hex color string (e.g. #FFFFFF) to an RGBColor object."""
        try:
            hex_cleaned = hex_str.lstrip("#")
            if len(hex_cleaned) == 3:
                hex_cleaned = "".join(c * 2 for c in hex_cleaned)
            r = int(hex_cleaned[0:2], 16)
            g = int(hex_cleaned[2:4], 16)
            b = int(hex_cleaned[4:6], 16)
            return RGBColor(r, g, b)
        except Exception:
            # Fallback to white on any parsing failure
            return RGBColor(255, 255, 255)
            
    def _get_alignment(self, align_str: str) -> PP_ALIGN:
        """Maps template alignment string to python-pptx PP_ALIGN enumeration."""
        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT
        }
        return align_map.get(align_str.lower(), PP_ALIGN.CENTER)

    def compile_ppt(self, parsed_refs: list[dict], template, output_path: str, validator):
        """
        Compiles the list of parsed/validated references into a PowerPoint presentation.
        """
        prs = Presentation()
        # Set 16:9 widescreen dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_layout = prs.slide_layouts[6] # Blank slide layout
        
        # Load book details for formatting references
        books = {b["id"]: b for b in self.repo.get_books()}
        
        # 1. Fetch all verses grouped by slides
        slides_data = []
        
        for ref in parsed_refs:
            val_res = validator.validate_reference(ref)
            if not val_res["valid"]:
                continue
                
            book_id = val_res["book_id"]
            book_meta = books.get(book_id)
            if not book_meta:
                continue
                
            verses = self.repo.get_verses(
                book_id=book_id,
                chapter=ref["chapter"],
                start_verse=ref["start_verse"],
                end_verse=ref["end_verse"]
            )
            
            # Format reference tag
            book_display_ko = book_meta["name_ko"]
            book_display_en = book_meta["name_en"]
            ch = ref["chapter"]
            v_start = ref["start_verse"]
            v_end = ref["end_verse"]
            v_range = f"{v_start}" if v_start == v_end else f"{v_start}-{v_end}"
            
            ref_ko = f"{book_display_ko} {ch}:{v_range}"
            ref_en = f"{book_display_en} {ch}:{v_range}"
            if template.ref_order == "ko_en":
                ref_label = f"{ref_ko} - {ref_en}"
            else:
                ref_label = f"{ref_en} - {ref_ko}"
            
            if template.split_verses:
                # One slide per verse
                for v in verses:
                    single_v_range = f"{v['verse']}"
                    single_ref_ko = f"{book_display_ko} {ch}:{single_v_range}"
                    single_ref_en = f"{book_display_en} {ch}:{single_v_range}"
                    if template.ref_order == "ko_en":
                        single_ref_label = f"{single_ref_ko} - {single_ref_en}"
                    else:
                        single_ref_label = f"{single_ref_en} - {single_ref_ko}"
                    
                    slides_data.append({
                        "ref_label": single_ref_label,
                        "verses": [v]
                    })
            else:
                # Combine verse range on one slide
                slides_data.append({
                    "ref_label": ref_label,
                    "verses": verses
                })
                
        # 2. Generate Slides
        bg_rgb = self._hex_to_rgb(template.bg_color)
        font_rgb = self._hex_to_rgb(template.font_color)
        alignment = self._get_alignment(template.alignment)
        
        for slide_info in slides_data:
            slide = prs.slides.add_slide(blank_layout)
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_rgb
            
            # Define text box bounds
            # Left: 0.8", Width: 11.733" (giving a margins of 0.8" on both sides)
            left = Inches(0.8)
            width = Inches(11.733)
            height = Inches(5.5)
            
            # Calculate top based on verse location setting
            if template.verse_location == "top":
                top = Inches(0.6)
            elif template.verse_location == "bottom":
                top = Inches(1.8)
            else: # center
                top = Inches(1.0)
                
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # Formulate text items
            verses = slide_info["verses"]
            
            # Build text strings based on configuration
            ko_texts = []
            en_texts = []
            
            for v in verses:
                prefix = f"{v['verse']} " if len(verses) > 1 else ""
                
                # Check for None/NULL translation values
                ngayok_val = v.get("ngayok")
                ngayok_txt = ngayok_val if ngayok_val else ""
                
                niv_val = v.get("niv")
                niv_txt = niv_val if niv_val else ""
                
                ko_texts.append(f"{prefix}{ngayok_txt}")
                en_texts.append(f"{prefix}{niv_txt}")
                
            ko_combined = " ".join(ko_texts)
            en_combined = " ".join(en_texts)
            
            # Populate paragraphs inside the text frame
            is_first = True
            
            def add_ref_para():
                nonlocal is_first
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p.text = slide_info["ref_label"]
                p.font.name = template.font_en
                p.font.size = Pt(int(template.font_size_en * 0.8)) # Slightly smaller reference font size
                p.font.color.rgb = font_rgb
                p.alignment = alignment
                p.space_after = Pt(14)
                
            def add_ko_para():
                nonlocal is_first
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p.text = ko_combined
                p.font.name = template.font_ko
                p.font.size = Pt(template.font_size_ko)
                p.font.color.rgb = font_rgb
                p.alignment = alignment
                p.space_after = Pt(16)
                p.line_spacing = template.line_spacing
                
            def add_en_para():
                nonlocal is_first
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p.text = en_combined
                p.font.name = template.font_en
                p.font.size = Pt(template.font_size_en)
                p.font.color.rgb = font_rgb
                p.alignment = alignment
                p.space_after = Pt(16)
                p.line_spacing = template.line_spacing

            # Execute ordering based on settings
            if template.ref_location == "top":
                add_ref_para()
                
            if template.bilingual_order == "ko_en":
                add_ko_para()
                add_en_para()
            elif template.bilingual_order == "en_ko":
                add_en_para()
                add_ko_para()
            elif template.bilingual_order == "ko_only":
                add_ko_para()
            elif template.bilingual_order == "en_only":
                add_en_para()
                
            if template.ref_location == "bottom":
                add_ref_para()
                
        # Save presentation
        prs.save(output_path)
